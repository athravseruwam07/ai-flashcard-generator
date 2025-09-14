# utils/parsing.py
# file reading + basic cleanup helpers

from __future__ import annotations

import re
from typing import Any

from pypdf import PdfReader
from docx import Document
from pptx import Presentation  # new: pptx support


def read_file(uploaded_file: Any) -> str:
    """auto-detect by extension and extract text; raise helpful error on failure."""
    name = getattr(uploaded_file, "name", "uploaded")
    lower = name.lower()

    if lower.endswith(".pdf"):
        try:
            reader = PdfReader(uploaded_file)
            pages = []
            for p in reader.pages:
                pages.append(p.extract_text() or "")
            text = "\n\n".join(pages)
            if not text.strip():
                raise ValueError("pdf text extraction returned empty text")
            return text
        except Exception as e:
            raise RuntimeError(f"pdf parsing failed: {e}")

    if lower.endswith(".docx"):
        try:
            doc = Document(uploaded_file)
            text = "\n".join([p.text for p in doc.paragraphs])
            if not text.strip():
                raise ValueError("docx appears to contain no readable paragraphs")
            return text
        except Exception as e:
            raise RuntimeError(f"docx parsing failed: {e}")

    if lower.endswith(".pptx"):
        try:
            prs = Presentation(uploaded_file)
            buf = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = shape.text or ""
                        if t.strip():
                            buf.append(t.strip())
            text = "\n\n".join(buf)
            if not text.strip():
                raise ValueError("pptx appears to contain no readable text")
            return text
        except Exception as e:
            raise RuntimeError(f"pptx parsing failed: {e}")

    if lower.endswith(".txt"):
        try:
            content = uploaded_file.read()
            if isinstance(content, bytes):
                try:
                    return content.decode("utf-8")
                except UnicodeDecodeError:
                    return content.decode("latin-1", errors="ignore")
            return str(content)
        except Exception as e:
            raise RuntimeError(f"txt read failed: {e}")

    raise ValueError("unsupported file type. use pdf, docx, pptx, or txt")


def clean_text(text: str) -> str:
    """normalize whitespace and remove obvious repeated headers/footers."""
    s = text.replace("\r\n", "\n").replace("\r", "\n")
    # collapse excessive spaces
    s = re.sub(r"\u00a0", " ", s)  # non-breaking space
    s = re.sub(r"[ \t]+", " ", s)
    # dedupe newlines to max 2
    s = re.sub(r"\n{3,}", "\n\n", s)

    # try to remove super-common short lines (likely headers/footers)
    lines = [ln.strip() for ln in s.split("\n")]
    freq = {}
    for ln in lines:
        if 0 < len(ln) <= 60:
            freq[ln] = freq.get(ln, 0) + 1
    repeated = {k for k, v in freq.items() if v >= 3}
    if repeated:
        lines = [ln for ln in lines if ln not in repeated]
    s = "\n".join(lines)

    return s.strip()